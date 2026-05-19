from __future__ import annotations

import re
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any
from urllib.parse import unquote
from zipfile import BadZipFile, ZipFile

from pptx import Presentation

from .schemas import SlideText


SLIDE_XML_RE = re.compile(r"^ppt/slides/slide(\d+)\.xml$")
REL_NS = "{http://schemas.openxmlformats.org/package/2006/relationships}"


def extract_pptx_slides(path: Path) -> list[SlideText]:
    zip_slides = _extract_slides_from_zip(path)

    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        if zip_slides:
            return zip_slides
        raise ValueError("Не удалось прочитать pptx-файл") from exc

    slides: list[SlideText] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        seen: set[str] = set()

        for shape in slide.shapes:
            _collect_shape_text(shape, parts, seen)

        _collect_xml_text(slide.element, parts, seen)
        slides.append(SlideText(index=index, text="\n".join(parts)))

    if not slides:
        return zip_slides

    zip_text_by_index = {slide.index: slide.text for slide in zip_slides}
    merged_slides: list[SlideText] = []
    for slide in slides:
        text = slide.text or zip_text_by_index.get(slide.index, "")
        merged_slides.append(SlideText(index=slide.index, text=text))

    for slide in zip_slides:
        if slide.index > len(merged_slides):
            merged_slides.append(slide)

    return merged_slides


def _extract_slides_from_zip(path: Path) -> list[SlideText]:
    try:
        with ZipFile(path) as archive:
            slide_names = sorted(
                (
                    (int(match.group(1)), name)
                    for name in archive.namelist()
                    if (match := SLIDE_XML_RE.match(name))
                ),
                key=lambda item: item[0],
            )

            slides: list[SlideText] = []
            for index, name in slide_names:
                slides.append(SlideText(index=index, text=_extract_zip_slide_text(archive, name)))
            return slides
    except (BadZipFile, KeyError, ET.ParseError):
        return []


def _extract_zip_slide_text(archive: ZipFile, slide_name: str) -> str:
    parts: list[str] = []
    seen: set[str] = set()

    _append_xml_text(archive.read(slide_name), parts, seen)

    for related_name in _iter_slide_related_xml_names(archive, slide_name):
        try:
            _append_xml_text(archive.read(related_name), parts, seen)
        except (KeyError, ET.ParseError):
            continue

    return "\n".join(parts)


def _iter_slide_related_xml_names(archive: ZipFile, slide_name: str) -> list[str]:
    slide_path = Path(slide_name)
    rels_name = slide_path.parent / "_rels" / f"{slide_path.name}.rels"
    try:
        root = ET.fromstring(archive.read(str(rels_name).replace("\\", "/")))
    except (KeyError, ET.ParseError):
        return []

    names: list[str] = []
    archive_names = set(archive.namelist())
    for relationship in root.findall(f"{REL_NS}Relationship"):
        target = relationship.attrib.get("Target", "")
        if not target or relationship.attrib.get("TargetMode") == "External":
            continue

        normalized = _normalize_related_target(slide_path.parent, target)
        if normalized.endswith(".xml") and normalized in archive_names:
            names.append(normalized)

    return names


def _normalize_related_target(base_dir: Path, target: str) -> str:
    target = unquote(target).replace("\\", "/")
    if target.startswith("/"):
        return target.lstrip("/")

    parts: list[str] = []
    for part in f"{base_dir.as_posix()}/{target}".split("/"):
        if not part or part == ".":
            continue
        if part == "..":
            if parts:
                parts.pop()
            continue
        parts.append(part)

    return "/".join(parts)


def _extract_xml_text(xml: bytes) -> str:
    parts: list[str] = []
    seen: set[str] = set()
    _append_xml_text(xml, parts, seen)
    return "\n".join(parts)


def _append_xml_text(xml: bytes, parts: list[str], seen: set[str]) -> None:
    root = ET.fromstring(xml)

    for paragraph in root.iter():
        if _local_name(paragraph.tag) != "p":
            continue

        text = "".join(
            node.text or ""
            for node in paragraph.iter()
            if _local_name(node.tag) == "t"
        )
        _append_text(text, parts, seen)

    for node in root.iter():
        if _local_name(node.tag) in {"t", "v"}:
            _append_text(node.text, parts, seen)


def _collect_shape_text(shape: Any, parts: list[str], seen: set[str]) -> None:
    if getattr(shape, "has_text_frame", False) and shape.text_frame:
        for paragraph in shape.text_frame.paragraphs:
            _append_text(paragraph.text, parts, seen)

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    _append_text(paragraph.text, parts, seen)

    if hasattr(shape, "shapes"):
        for nested_shape in shape.shapes:
            _collect_shape_text(nested_shape, parts, seen)


def _collect_xml_text(element: Any, parts: list[str], seen: set[str]) -> None:
    for paragraph in element.iter():
        if _local_name(paragraph.tag) != "p":
            continue

        text = "".join(
            node.text or ""
            for node in paragraph.iter()
            if _local_name(node.tag) == "t"
        )
        _append_text(text, parts, seen)


def _append_text(text: str | None, parts: list[str], seen: set[str]) -> None:
    normalized = " ".join((text or "").split())
    if normalized and normalized not in seen:
        parts.append(normalized)
        seen.add(normalized)


def _local_name(tag: Any) -> str:
    if not isinstance(tag, str):
        return ""
    return tag.rsplit("}", 1)[-1]
